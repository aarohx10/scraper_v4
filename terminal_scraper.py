import asyncio
import json
import re
import requests
from urllib.parse import urlparse, urljoin, quote_plus
from bs4 import BeautifulSoup
import time
import random
import sys
import socket
import os
import mimetypes
import fitz  # PyMuPDF
import docx
import openpyxl
import pptx
import tldextract
from tqdm import tqdm
import aiohttp
import aiofiles
from asyncio import Semaphore
import logging

logging.basicConfig(level=logging.INFO)

def google_search_urls(query, num_results=20):
    """Scrape Google search results for a query using requests."""
    print(f"Searching the web for: {query}")
    
    # Extract company and website from query
    company_name = None
    website = None
    
    # Extract website if present
    website_match = re.search(r'www\.[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+', query)
    if website_match:
        website = website_match.group(0)
        # Try to extract company name from parts before the website
        parts = query.split(website)
        if parts[0].strip():
            company_name = parts[0].strip().rstrip(',')
    else:
        # If no website, assume entire query is company name
        company_name = query
        
    print(f"Detected company: {company_name}")
    print(f"Detected website: {website}")
    
    urls = []
    
    # 1. Try the website directly if provided
    if website:
        if not website.startswith('http'):
            website_url = f"https://{website}"
        else:
            website_url = website
            
        urls.append(website_url)
        
        # Add common subpages
        for path in ['', '/about', '/about-us', '/company', '/team', '/contact', '/products', 
                     '/services', '/blog', '/news', '/careers', '/jobs']:
            if path == '' and website_url in urls:
                continue
            urls.append(f"{website_url}{path}")
    
    # 2. Add LinkedIn search for the company
    if company_name:
        urls.append(f"https://www.linkedin.com/company/{company_name.lower().replace(' ', '-')}")
        urls.append(f"https://www.linkedin.com/search/results/companies/?keywords={quote_plus(company_name)}")
    
    # 3. Add other common places to find company info
    if company_name:
        for site in ['crunchbase.com', 'bloomberg.com', 'zoominfo.com', 'glassdoor.com', 
                     'hoovers.com', 'owler.com', 'indeed.com', 'angel.co']:
            urls.append(f"https://www.google.com/search?q={quote_plus(company_name)}+site:{site}")
    
    # 4. Search for news about the company
    if company_name:
        urls.append(f"https://www.google.com/search?q={quote_plus(company_name)}+news")
        urls.append(f"https://www.google.com/search?q={quote_plus(company_name)}+press+release")
    
    # 5. Check for social media presence
    if company_name:
        for social in ['twitter.com', 'facebook.com', 'instagram.com', 'youtube.com']:
            if website:
                company_handle = website.split('.')[1]
                urls.append(f"https://{social}/{company_handle}")
            else:
                company_handle = company_name.lower().replace(' ', '')
                urls.append(f"https://{social}/{company_handle}")
    
    # 6. Search for company reviews
    if company_name:
        urls.append(f"https://www.google.com/search?q={quote_plus(company_name)}+reviews")
        urls.append(f"https://www.google.com/search?q={quote_plus(company_name)}+ratings")
    
    # Remove duplicates and limit results
    urls = list(dict.fromkeys(urls))
    
    return urls[:num_results]

def is_valid_url(url):
    """Check if a URL is valid and reachable."""
    try:
        result = urlparse(url)
        return all([result.scheme, result.netloc])
    except:
        return False

def normalize_url(base_url, url):
    """Convert relative URLs to absolute URLs."""
    if is_valid_url(url):
        return url
    return urljoin(base_url, url)

def crawl_page(url):
    """Extract text from a webpage."""
    print(f"Crawling {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,image/apng,*/*;q=0.8",
            "Referer": "https://www.google.com/"
        }
        
        # Set a reasonable timeout
        response = requests.get(url, timeout=20, headers=headers, allow_redirects=True)
        response.raise_for_status()
        
        content_type = response.headers.get('Content-Type', '').lower()
        if not ('text/html' in content_type or 'application/xhtml+xml' in content_type):
            return f"[Not HTML content: {content_type}]"
        
        # Parse HTML content
        soup = BeautifulSoup(response.text, "html.parser")
        
        # Extract title
        title = soup.title.string if soup.title else "No title"
        
        # Extract meta description
        meta_desc = ""
        meta_tag = soup.find("meta", attrs={"name": "description"})
        if meta_tag and meta_tag.get("content"):
            meta_desc = meta_tag["content"]
            
        # Extract text content
        for script in soup(["script", "style"]):
            script.extract()
            
        # Try to get structured content first
        structured_content = ""
        
        # Check for specific content sections
        article = soup.find(['article', 'main'])
        if article:
            structured_content = article.get_text(separator="\n", strip=True)
        
        # If no article/main, look for content divs
        if not structured_content:
            content_divs = soup.find_all(['div', 'section'], class_=lambda c: c and any(x in str(c).lower() for x in ['content', 'main', 'article', 'body']))
            if content_divs:
                for div in content_divs:
                    structured_content += div.get_text(separator="\n", strip=True) + "\n\n"
        
        # If still no content, get all paragraphs
        if not structured_content:
            paragraphs = soup.find_all('p')
            structured_content = "\n\n".join([p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)])
        
        # If still nothing, get all text
        if not structured_content:
            structured_content = soup.get_text(separator="\n", strip=True)
        
        # Special handling for LinkedIn
        if 'linkedin.com' in url:
            linkedin_sections = []
            
            # Company overview
            overview = soup.find_all(['section', 'div'], class_=lambda c: c and 'overview' in str(c).lower())
            if overview:
                for section in overview:
                    linkedin_sections.append("COMPANY OVERVIEW:\n" + section.get_text(separator="\n", strip=True))
            
            # About section
            about = soup.find_all(['section', 'div'], class_=lambda c: c and 'about' in str(c).lower())
            if about:
                for section in about:
                    linkedin_sections.append("ABOUT SECTION:\n" + section.get_text(separator="\n", strip=True))
            
            # Employees/people
            people = soup.find_all(['section', 'div'], class_=lambda c: c and 'employee' in str(c).lower() or 'people' in str(c).lower())
            if people:
                for section in people:
                    linkedin_sections.append("EMPLOYEES/PEOPLE:\n" + section.get_text(separator="\n", strip=True))
                    
            if linkedin_sections:
                structured_content = "\n\n".join(linkedin_sections)
        
        # Special handling for review sites
        if any(site in url for site in ['glassdoor.com', 'indeed.com', 'yelp.com']):
            reviews = soup.find_all(['div', 'section'], class_=lambda c: c and 'review' in str(c).lower())
            if reviews:
                reviews_text = []
                for section in reviews:
                    reviews_text.append(section.get_text(separator="\n", strip=True))
                structured_content = "REVIEWS:\n\n" + "\n\n".join(reviews_text)
        
        # Format the final output
        full_content = f"SOURCE: {url}\nTITLE: {title}\nDESCRIPTION: {meta_desc}\n\nFULL CONTENT:\n{structured_content}"
        
        return full_content
        
    except requests.exceptions.RequestException as e:
        print(f"Request error crawling {url}: {e}")
        return f"[Could not retrieve content from {url}: {str(e)}]"
    except Exception as e:
        print(f"Error crawling {url}: {e}")
        return f"[Error processing {url}: {str(e)}]"

async def process_url(url):
    """Process a single URL and extract all relevant information."""
    # Add some random delay to avoid rate limiting
    await asyncio.sleep(random.uniform(0.5, 2.0))
    
    page_text = crawl_page(url)
    
    # Don't trim content for full detailed output
    return {
        "url": url,
        "content": page_text  # No character limit here
    }

def extract_emails(text):
    return re.findall(r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+", text)

def extract_phones(text):
    return re.findall(r"\+?\d[\d\s().-]{7,}\d", text)

def extract_links(soup, base_url):
    links = set()
    for a in soup.find_all('a', href=True):
        href = a['href']
        if href.startswith('http'):
            links.add(href)
        elif href.startswith('/'):
            links.add(urljoin(base_url, href))
    return links

def extract_documents(soup, base_url):
    doc_links = set()
    for a in soup.find_all('a', href=True):
        href = a['href']
        if any(href.lower().endswith(ext) for ext in ['.pdf', '.docx', '.xlsx', '.pptx', '.txt']):
            if href.startswith('http'):
                doc_links.add(href)
            elif href.startswith('/'):
                doc_links.add(urljoin(base_url, href))
    return doc_links

def extract_text_from_pdf(path):
    text = ""
    try:
        with fitz.open(path) as doc:
            for page in doc:
                text += page.get_text()
    except Exception as e:
        logging.error(f"PDF extraction error for {path}: {e}")
        text = f"[PDF extraction error: {e}]"
    return text

def extract_text_from_docx(path):
    try:
        doc = docx.Document(path)
        return '\n'.join([p.text for p in doc.paragraphs])
    except Exception as e:
        logging.error(f"DOCX extraction error for {path}: {e}")
        return f"[DOCX extraction error: {e}]"

def extract_text_from_xlsx(path):
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                text.append('\t'.join([str(cell) if cell is not None else '' for cell in row]))
        return '\n'.join(text)
    except Exception as e:
        logging.error(f"XLSX extraction error for {path}: {e}")
        return f"[XLSX extraction error: {e}]"

def extract_text_from_pptx(path):
    try:
        prs = pptx.Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        logging.error(f"PPTX extraction error for {path}: {e}")
        return f"[PPTX extraction error: {e}]"

def extract_text_from_txt(path):
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logging.error(f"TXT extraction error for {path}: {e}")
        return f"[TXT extraction error: {e}]"

def download_file(url, dest_folder="downloads"):
    os.makedirs(dest_folder, exist_ok=True)
    local_filename = os.path.join(dest_folder, os.path.basename(url.split('?')[0]))
    try:
        r = requests.get(url, stream=True, timeout=30)
        r.raise_for_status()
        with open(local_filename, 'wb') as f:
            for chunk in r.iter_content(chunk_size=8192):
                f.write(chunk)
        return local_filename
    except Exception as e:
        logging.error(f"Exception downloading {url}: {e}")
        return None

def extract_detailed_content(soup, url):
    content = {}
    content['url'] = url
    content['title'] = soup.title.string if soup.title else "No title"
    content['meta_description'] = ''
    meta_tag = soup.find("meta", attrs={"name": "description"})
    if meta_tag and meta_tag.get("content"):
        content['meta_description'] = meta_tag["content"]
    content['headings'] = {f"h{i}": [h.get_text(strip=True) for h in soup.find_all(f"h{i}")] for i in range(1,7)}
    content['paragraphs'] = [p.get_text(strip=True) for p in soup.find_all('p') if p.get_text(strip=True)]
    content['tables'] = []
    for table in soup.find_all('table'):
        rows = []
        for tr in table.find_all('tr'):
            cells = [td.get_text(strip=True) for td in tr.find_all(['td','th'])]
            if cells:
                rows.append(cells)
        if rows:
            content['tables'].append(rows)
    content['lists'] = []
    for ul in soup.find_all(['ul','ol']):
        items = [li.get_text(strip=True) for li in ul.find_all('li')]
        if items:
            content['lists'].append(items)
    content['images'] = [{'src': img.get('src'), 'alt': img.get('alt','')} for img in soup.find_all('img')]
    page_text = soup.get_text(separator="\n", strip=True)
    content['emails'] = extract_emails(page_text)
    content['phones'] = extract_phones(page_text)
    content['links'] = list(extract_links(soup, url))
    return content

def recursive_crawl(start_url, max_pages=30, max_depth=2):
    visited = set()
    to_visit = [(start_url, 0)]
    domain = tldextract.extract(start_url).registered_domain
    all_content = []
    with tqdm(total=max_pages, desc=f"Crawling {start_url}") as pbar:
        while to_visit and len(visited) < max_pages:
            url, depth = to_visit.pop(0)
            if url in visited or depth > max_depth:
                continue
            try:
                response = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
                if response.status_code != 200:
                    continue
                soup = BeautifulSoup(response.text, "lxml")
                content = extract_detailed_content(soup, url)
                all_content.append(content)
                visited.add(url)
                pbar.update(1)
                # Find new links
                for link in extract_links(soup, url):
                    link_domain = tldextract.extract(link).registered_domain
                    if link_domain == domain and link not in visited:
                        to_visit.append((link, depth+1))
            except Exception as e:
                continue
    return all_content

def process_documents(doc_links):
    docs = []
    for url in tqdm(doc_links, desc="Downloading documents"):
        local_path = download_file(url)
        if not local_path:
            continue
        ext = os.path.splitext(local_path)[1].lower()
        if ext == '.pdf':
            text = extract_text_from_pdf(local_path)
        elif ext == '.docx':
            text = extract_text_from_docx(local_path)
        elif ext == '.xlsx':
            text = extract_text_from_xlsx(local_path)
        elif ext == '.pptx':
            text = extract_text_from_pptx(local_path)
        elif ext == '.txt':
            text = extract_text_from_txt(local_path)
        else:
            text = '[Unknown document type]'
        docs.append({'url': url, 'local_path': local_path, 'content': text})
    return docs

CONCURRENT_REQUESTS = 10

async def fetch_page(session, url, sem):
    async with sem:
        try:
            headers = {"User-Agent": "Mozilla/5.0"}
            async with session.get(url, timeout=15, headers=headers) as resp:
                if resp.status != 200:
                    logging.warning(f"Failed to fetch {url}: status {resp.status}")
                    return None
                text = await resp.text()
                return text
        except Exception as e:
            logging.error(f"Exception fetching {url}: {e}")
            return None

async def extract_detailed_content_async(session, url, sem):
    html = await fetch_page(session, url, sem)
    if not html:
        return None
    soup = BeautifulSoup(html, "lxml")
    return extract_detailed_content(soup, url)

async def recursive_crawl_async(start_url, max_pages=30, max_depth=2):
    visited = set()
    to_visit = [(start_url, 0)]
    domain = tldextract.extract(start_url).registered_domain
    all_content = []
    sem = Semaphore(CONCURRENT_REQUESTS)
    async with aiohttp.ClientSession() as session:
        while to_visit and len(visited) < max_pages:
            batch = []
            while to_visit and len(batch) < CONCURRENT_REQUESTS and len(visited) + len(batch) < max_pages:
                url, depth = to_visit.pop(0)
                if url in visited or depth > max_depth:
                    continue
                batch.append((url, depth))
            tasks = [extract_detailed_content_async(session, url, sem) for url, _ in batch]
            results = await asyncio.gather(*tasks)
            for (url, depth), content in zip(batch, results):
                if content:
                    all_content.append(content)
                    visited.add(url)
                    # Find new links
                    soup = BeautifulSoup(await fetch_page(session, url, sem), "lxml")
                    for link in extract_links(soup, url):
                        link_domain = tldextract.extract(link).registered_domain
                        if link_domain == domain and link not in visited:
                            to_visit.append((link, depth+1))
    return all_content

async def download_file_async(session, url, dest_folder="downloads", sem=None):
    os.makedirs(dest_folder, exist_ok=True)
    local_filename = os.path.join(dest_folder, os.path.basename(url.split('?')[0]))
    try:
        async with sem:
            async with session.get(url, timeout=30) as resp:
                if resp.status != 200:
                    logging.warning(f"Failed to download {url}: status {resp.status}")
                    return None
                async with aiofiles.open(local_filename, 'wb') as f:
                    async for chunk in resp.content.iter_chunked(8192):
                        await f.write(chunk)
        return local_filename
    except Exception as e:
        logging.error(f"Exception downloading {url}: {e}")
        return None

async def process_documents_async(doc_links):
    docs = []
    sem = Semaphore(CONCURRENT_REQUESTS)
    async with aiohttp.ClientSession() as session:
        tasks = [download_file_async(session, url, sem=sem) for url in doc_links]
        local_paths = await asyncio.gather(*tasks)
    for url, local_path in zip(doc_links, local_paths):
        if not local_path:
            continue
        ext = os.path.splitext(local_path)[1].lower()
        if ext == '.pdf':
            text = extract_text_from_pdf(local_path)
        elif ext == '.docx':
            text = extract_text_from_docx(local_path)
        elif ext == '.xlsx':
            text = extract_text_from_xlsx(local_path)
        elif ext == '.pptx':
            text = extract_text_from_pptx(local_path)
        elif ext == '.txt':
            text = extract_text_from_txt(local_path)
        else:
            text = '[Unknown document type]'
        docs.append({'url': url, 'local_path': local_path, 'content': text})
    return docs

async def main(query):
    print(f"Starting comprehensive company research for: {query}")
    urls = google_search_urls(query, num_results=10)
    all_results = []
    for url in urls:
        print(f"\nRecursively crawling: {url}")
        site_content = await recursive_crawl_async(url, max_pages=30, max_depth=2)
        # Collect all document links from all pages
        doc_links = set()
        for page in site_content:
            html = page.get('html') if 'html' in page else None
            if not html:
                try:
                    async with aiohttp.ClientSession() as session:
                        html = await fetch_page(session, page['url'], Semaphore(1))
                except Exception:
                    continue
            soup = BeautifulSoup(html, "lxml")
            doc_links.update(extract_documents(soup, page['url']))
        documents = await process_documents_async(list(doc_links))
        all_results.append({'site': url, 'pages': site_content, 'documents': documents})
    print("\nResearch completed. Returning all detailed data.")
    return all_results

if __name__ == "__main__":
    # If arguments provided, use them as query
    if len(sys.argv) > 1:
        query = " ".join(sys.argv[1:])
    else:
        # Otherwise prompt for input
        query = input("Enter company name and/or website: ")
    
    if not query:
        print("No query provided. Exiting.")
        sys.exit(1)
    
    print(f"Starting comprehensive research for: {query}")
    asyncio.run(main(query))